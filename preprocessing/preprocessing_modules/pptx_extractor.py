from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any
from PIL import Image
from io import BytesIO
import requests
from concurrent.futures import ThreadPoolExecutor, as_completed
import tempfile
import os


# This part requires a config file or environment variable for the API key.
# For demonstration, you can replace OCR_SPACE_API_KEY with your actual key string.
# from config.config import OCR_SPACE_API_KEY
from config.config import OCR_SPACE_API_KEY # <-- IMPORTANT: Replace with your actual OCR Space API key
from logger.custom_logger import CustomLogger

# module logger
logger = CustomLogger().get_logger(__file__)

API_URL = "https://api.ocr.space/parse/image"

assert OCR_SPACE_API_KEY != "YOUR_API_KEY_HERE", "OCR_SPACE_API_KEY not set."

def ocr_space_file(filename, api_key=OCR_SPACE_API_KEY, overlay=False, language="eng"):
    """Extract text from image file using OCR Space API"""
    payload = {
        "isOverlayRequired": overlay,
        "apikey": api_key,
        "language": language,
        "detectOrientation": True,
        "scale": True,
        "isTable": False,
        "OCREngine": 2
    }
    try:
        with open(filename, "rb") as f:
            response = requests.post(API_URL, files={filename: f}, data=payload, timeout=30)
        
        if response.status_code != 200:
            return filename, f"API Error: HTTP {response.status_code}"
            
        parsed = response.json()
        
        if parsed.get("OCRExitCode") == 1:
            parsed_text = parsed.get("ParsedResults", [{}])[0].get("ParsedText", "")
            return filename, parsed_text
        else:
            error_msg = parsed.get("ErrorMessage", ["Unknown error"])[0] if parsed.get("ErrorMessage") else "Unknown OCR error"
            return filename, f"OCR Error: {error_msg}"
            
    except requests.exceptions.Timeout:
        return filename, "Error: Request timeout"
    except requests.exceptions.RequestException as e:
        return filename, f"Error: Network error - {str(e)}"
    except Exception as e:
        return filename, f"Error: {e}"

def batch_ocr_parallel(filenames, max_workers=5):
    """Process multiple image files in parallel using OCR Space API"""
    results = {}
    if not filenames:
        return results
        
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_file = {executor.submit(ocr_space_file, fname): fname for fname in filenames}
        for future in as_completed(future_to_file):
            fname, text = future.result()
            results[fname] = text
    return results

def extract_pptx_with_meta(pptx_path: str) -> List[Dict[str, Any]]:
    """Extract content from PPTX with metadata, using OCR Space API for images"""
    prs = Presentation(pptx_path)
    all_slides_content = []
    
    # First pass: extract all images and save them temporarily
    temp_image_files = []
    image_to_shape_mapping = {}
    
    temp_dir = tempfile.mkdtemp()
    
    try:
        # Extract all images first
        logger.info("Extracting images from PPTX to temporary directory", temp_dir=temp_dir)
        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = Image.open(BytesIO(shape.image.blob))
                        temp_file = os.path.join(temp_dir, f"slide_{slide_index}shape{shape_index}.png")
                        img.save(temp_file, 'PNG')
                        temp_image_files.append(temp_file)
                        image_to_shape_mapping[temp_file] = (slide_index, shape_index)
                        logger.info("Extracted image from slide", slide=slide_index, shape=shape_index, file=temp_file)
                    except Exception as e:
                        logger.error("Failed to extract image from slide", slide=slide_index, shape=shape_index, error=str(e))
        
        # Process all images in parallel using OCR Space API
        logger.info("Processing images with OCR Space API", total_images=len(temp_image_files))
        ocr_results = batch_ocr_parallel(temp_image_files, max_workers=5)
        logger.info("OCR processing completed", processed=len(ocr_results))
        
        # Second pass: build the content structure
        for slide_index, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": slide_index + 1,
                "content_blocks": []
            }

            for shape_index, shape in enumerate(slide.shapes):
                content_block = {}

                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        content_block["type"] = "text"
                        content_block["content"] = text

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Find the corresponding OCR result
                    temp_file_key = None
                    for temp_file, (s_idx, sh_idx) in image_to_shape_mapping.items():
                        if s_idx == slide_index and sh_idx == shape_index:
                            temp_file_key = temp_file
                            break

                    # If we found a matching temp file, look up OCR results
                    if temp_file_key and temp_file_key in ocr_results:
                        ocr_text = ocr_results[temp_file_key].strip()
                        if ocr_text and not ocr_text.startswith("Error:"):
                            content_block["type"] = "image"
                            content_block["content"] = ocr_text
                            logger.info("OCR extracted from slide", slide=slide_index, preview=ocr_text[:100])
                        else:
                            content_block["type"] = "image"
                            content_block["content"] = f"[OCR failed: {ocr_text}]"
                    else:
                        content_block["type"] = "image"
                        content_block["content"] = "[OCR processing failed - no result found]"

                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        table = shape.table
                        content_block["type"] = "table"
                        table_content = "---Table---\n"
                        for row in table.rows:
                            row_content = ", ".join([cell.text.strip() for cell in row.cells])
                            table_content += row_content + "\n"
                        table_content += "-" * 11
                        content_block["content"] = table_content
                    except Exception as e:
                        content_block["type"] = "table"
                        content_block["content"] = f"[Table extraction failed: {str(e)}]"

                if content_block:
                    slide_data["content_blocks"].append(content_block)

            # Handle slide notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_data["content_blocks"].append({
                        "type": "notes",
                        "content": notes
                    })

            all_slides_content.append(slide_data)
        
    finally:
        # Clean up temporary files
        logger.info("Cleaning up temporary files", count=len(temp_image_files))
        for temp_file in temp_image_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                logger.error("Failed to remove temp file", file=temp_file, error=str(e))
        
        # Remove temp directory
        try:
            if os.path.exists(temp_dir):
                os.rmdir(temp_dir)
            logger.info("Temporary directory cleanup completed")
        except Exception as e:
            logger.error("Failed to remove temp directory", dir=temp_dir, error=str(e))
    
    return all_slides_content

# MODIFIED FUNCTION
def extract_pptx(filepath: str) -> List[Dict[str, Any]]:
    """
    Extracts content from a PPTX file and organizes it by slide number.

    Each slide's text, table data, and image OCR results are combined into a
    single string.

    Args:
        filepath: The path to the .pptx file.

    Returns:
        A list of dictionaries, where each dictionary represents a slide
        and has the format: {'page_num': slide_number, 'content': '...'}
    """
    slides_output = []
    
    # Get the detailed, structured data from the meta function
    structured_slides = extract_pptx_with_meta(filepath)
    
    # Iterate over each slide's data
    for slide_data in structured_slides:
        # Collect all content from the current slide's blocks into a list
        slide_content_blocks = []
        for block in slide_data["content_blocks"]:
            content = block.get("content", "").strip()
            if content:
                slide_content_blocks.append(content)
        
        # Join the content blocks for this slide into a single string
        full_slide_content = "\n".join(slide_content_blocks)
        
        # Create the dictionary in the desired format and add it to our results list
        # Only add slides that have extracted content
        if full_slide_content:
            slides_output.append({
                "page_num": slide_data["slide_number"],
                "content": full_slide_content
            })
            
    return slides_output

# Example usage:
if __name__ == '__main__':
    # Create a dummy presentation for testing
    from pptx.util import Inches

    # Create a dummy pptx file
    prs = Presentation()
    # Slide 1: Title and text
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide1.shapes.title
    title.text = "Slide 1: Title"
    content_box = slide1.shapes.placeholders[1]
    content_box.text = "This is some text on the first slide."
    
    # Slide 2: Table
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide2.shapes.title
    title.text = "Slide 2: Table"
    table = slide2.shapes.add_table(rows=2, cols=2, left=Inches(2), top=Inches(2), width=Inches(4), height=Inches(1.5)).table
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Cell A"
    table.cell(1, 1).text = "Cell B"
    
    # Slide 3: Blank (will be skipped in output)
    prs.slides.add_slide(prs.slide_layouts[6])

    # Slide 4: Notes
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    title.text = "Slide 4: With Notes"
    notes_slide = slide4.has_notes_slide and slide4.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This is a presenter note."

    dummy_pptx_path = "dummy_presentation.pptx"
    prs.save(dummy_pptx_path)

    print(f"Created a dummy presentation: {dummy_pptx_path}")
    print("-" * 20)

    # Run the extraction function
    # Note: OCR will not run as we haven't added images.
    # To test OCR, you would need to add an image to a slide.
    extracted_data = extract_pptx(dummy_pptx_path)
    
    # Print the result in a readable format
    import json
    print("Extraction Result:")
    print(json.dumps(extracted_data, indent=2))

    # Clean up the dummy file
    os.remove(dummy_pptx_path)
    print(f"\nRemoved dummy file: {dummy_pptx_path}")